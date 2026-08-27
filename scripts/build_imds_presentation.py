#!/usr/bin/env python3
"""Build the C-suite IMDS agentic-workflow briefing (storyboard + 8 beats)."""

from __future__ import annotations

import argparse
import os
import shutil
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


W = Inches(13.333)
H = Inches(7.5)

NAVY = RGBColor(0x0B, 0x1F, 0x3A)
NAVY_MID = RGBColor(0x14, 0x36, 0x58)
TEAL = RGBColor(0x0F, 0x6E, 0x6E)
GOLD = RGBColor(0xC4, 0xA3, 0x5A)
COPPER = RGBColor(0xC9, 0x6A, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF = RGBColor(0xF5, 0xF7, 0xF9)
SLATE = RGBColor(0x3A, 0x47, 0x58)
MUTED = RGBColor(0x6A, 0x75, 0x84)
DARK = RGBColor(0x1A, 0x23, 0x32)
RED = RGBColor(0xB4, 0x3A, 0x3A)
GREEN = RGBColor(0x1F, 0x8A, 0x72)
BLUE = RGBColor(0x2C, 0x5F, 0x8A)
LIGHT_GOLD = RGBColor(0xF8, 0xF1, 0xE0)
LIGHT_TEAL = RGBColor(0xE5, 0xF3, 0xF1)
LIGHT_RED = RGBColor(0xF8, 0xEB, 0xEB)
LINE = RGBColor(0xDE, 0xE3, 0xE9)

FONT = "Calibri"
ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "presentations" / "IMDS_Agentic_Workflow.pptx"
DOCS_PPTX = ROOT / "docs" / "IMDS_Agentic_Workflow.pptx"
DOCS_PDF = ROOT / "docs" / "IMDS_Agentic_Workflow.pdf"
SLIDES_DIR = ROOT / "docs" / "slides"
TOTAL = 12

# Military / ambiguous jargon that must never appear in this deck.
FORBIDDEN_JARGON = (
    "air cover",
    "cover fire",
    "battlespace",
    "war room",
    "kill chain",
    "kill switch",
    "kill-switch",
    "aircover",
)

# Narrative spine — presenter run sheet. Every content slide follows these beats.
STORYBOARD = [
    {
        "n": "1",
        "beat": "Executive Opening",
        "intent": "Title with project name, presenter, and date. Hook first.",
        "hook": "Fifty MDS land every day. Five thousand sit open. IMDS is a production gate — not a filing cabinet.",
        "goal": "Capture attention and establish urgency.",
        "time": "0:00–1:00",
    },
    {
        "n": "2",
        "beat": "Pain Points",
        "intent": "Visual of inefficiencies, missed same-day accept, rising launch cost.",
        "hook": "Inaction is already on the P&L: specialist hours, PPAP delay, OEM scorecard exposure.",
        "goal": "Make the audience feel the cost of inaction.",
        "time": "1:00–4:00",
    },
    {
        "n": "3",
        "beat": "Proposed Solution",
        "intent": "Diagram: ingest received MDS → orchestrate Check → PASS/FAIL output.",
        "hook": "The live agent already runs our IMDS account: PASS accept-forward-propose; FAIL reject.",
        "goal": "Position AI as a strategic enabler, not just a technical tool.",
        "time": "4:00–8:00",
    },
    {
        "n": "4",
        "beat": "Business Impact",
        "intent": "Before/after: manual inbox versus the agentic workflow.",
        "hook": "Hours return to specialists. MDS stop stalling PPAP. Decisions are logged.",
        "goal": "Show tangible transformation.",
        "time": "8:00–10:30",
    },
    {
        "n": "5",
        "beat": "Budget & ROI",
        "intent": "Implementation, training, maintenance. Payback window.",
        "hook": "Internal effort, not a software RFP. Pilot hours pay back inside one quarter.",
        "goal": "Demonstrate financial viability and risk mitigation.",
        "time": "10:30–13:30",
    },
    {
        "n": "6",
        "beat": "Implementation Roadmap",
        "intent": "Timeline: Pilot → Scale → Optimize, with governance.",
        "hook": "Controlled rollout: 20 MDS, logged decisions, executive sponsorship, emergency halt.",
        "goal": "Build confidence in execution.",
        "time": "13:30–15:30",
    },
    {
        "n": "7",
        "beat": "Case Studies",
        "intent": "External automotive / agentic-automation proof, labeled illustrative.",
        "hook": "Peers already automated IMDS-class work. We keep OEM knowledge in-house.",
        "goal": "Provide credibility and reduce perceived risk.",
        "time": "15:30–17:30",
    },
    {
        "n": "8",
        "beat": "Call to Action",
        "intent": "Bold line: Invest in Agentic AI Today. Specific ask.",
        "hook": "Approve the 20-MDS pilot and the internal budget line this week.",
        "goal": "Drive decision and secure buy-in.",
        "time": "17:30–20:00",
    },
]


def _set_run(run, text, size, color, bold=False, italic=False):
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.east_asian_font = FONT


def _fill_line(shape, fill, line=None, width=Pt(1)):
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = width


def add_text(slide, l, t, w, h, lines):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, item in enumerate(lines):
        text, size, color = item[0], item[1], item[2]
        bold = item[3] if len(item) > 3 else False
        align = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        italic = item[5] if len(item) > 5 else False
        after = item[6] if len(item) > 6 else 6
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(after)
        if p.runs:
            _set_run(p.runs[0], text, size, color, bold, italic)
            for extra in p.runs[1:]:
                extra.text = ""
        else:
            r = p.add_run()
            _set_run(r, text, size, color, bold, italic)
    return box


def add_bullets(slide, l, t, w, h, items, size=14, color=DARK, spacing=8):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        p.space_after = Pt(spacing)
        r = p.add_run() if not p.runs else p.runs[0]
        _set_run(r, "▸  " + item, size, color, False)
    return box


def rect(slide, l, t, w, h, fill, line=None, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE,
        l, t, w, h,
    )
    _fill_line(shape, fill, line)
    if radius is not None:
        try:
            shape.adjustments[0] = radius
        except Exception:
            pass
    return shape


def chevron(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, l, t, w, h)
    _fill_line(s, fill)
    return s


def oval(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    _fill_line(s, fill)
    return s


def card(slide, l, t, w, h, fill=WHITE, line=LINE):
    return rect(slide, l, t, w, h, fill, line, radius=0.06)


def footer(slide, n):
    rect(slide, Inches(0), Inches(7.28), W, Inches(0.22), OFF)
    rect(slide, Inches(0), Inches(7.28), W, Emu(12700), GOLD)
    add_text(
        slide,
        Inches(0.45),
        Inches(7.30),
        Inches(9.8),
        Inches(0.18),
        [("INTERNAL CONFIDENTIAL  ·  Johnson Electric  ·  Supplier Quality  ·  27 August 2026", 9, RGBColor(0x8A, 0x96, 0xA4), False)],
    )
    add_text(
        slide,
        Inches(11.2),
        Inches(7.30),
        Inches(1.7),
        Inches(0.18),
        [(f"{n}  /  {TOTAL}", 9, RGBColor(0x8A, 0x96, 0xA4), False, PP_ALIGN.RIGHT)],
    )


def header(slide, kicker, title, subtitle=None):
    hh = Inches(1.16) if subtitle else Inches(0.92)
    rect(slide, Inches(0), Inches(0), W, hh, NAVY)
    rect(slide, Inches(0), hh, W, Emu(14000), GOLD)
    add_text(slide, Inches(0.5), Inches(0.10), Inches(12.3), Inches(0.22), [(kicker.upper(), 10, GOLD, True)])
    add_text(slide, Inches(0.5), Inches(0.32), Inches(12.3), Inches(0.40), [(title, 22, WHITE, True)])
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.74), Inches(12.3), Inches(0.36), [(subtitle, 12, RGBColor(0xC5, 0xD0, 0xDC), False)])
    return hh


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def kpi(slide, l, t, w, h, value, label, sub, fill=NAVY):
    card(slide, l, t, w, h, fill, None)
    add_text(slide, l + Inches(0.16), t + Inches(0.12), w - Inches(0.28), Inches(0.42), [(value, 26, WHITE, True)])
    add_text(slide, l + Inches(0.16), t + Inches(0.54), w - Inches(0.28), Inches(0.36), [(label, 12, GOLD, False)])
    add_text(slide, l + Inches(0.16), t + Inches(0.90), w - Inches(0.28), Inches(0.42), [(sub, 11, RGBColor(0xC5, 0xD0, 0xDC), False)])


def set_cell(cell, text, size=11, color=DARK, bold=False, fill=None, align=PP_ALIGN.LEFT):
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    else:
        cell.fill.background()
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(2)
    p.space_after = Pt(2)
    if p.runs:
        _set_run(p.runs[0], text, size, color, bold)
        for extra in p.runs[1:]:
            extra.text = ""
    else:
        r = p.add_run()
        _set_run(r, text, size, color, bold)


def _speech_goal(beat_n: str, extra: str) -> str:
    row = next(b for b in STORYBOARD if b["n"] == beat_n)
    return (
        f"BEAT {row['n']} · {row['beat']}  ({row['time']})\n"
        f"SLIDE INTENT: {row['intent']}\n"
        f"SPOKEN HOOK: {row['hook']}\n"
        f"GOAL: {row['goal']}\n\n"
        f"{extra}"
    )


# --- Slides ----------------------------------------------------------------
def s01_title(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, NAVY)
    rect(sl, 0, 0, Inches(0.18), H, GOLD)
    rect(sl, Inches(0.18), 0, Inches(0.08), H, TEAL)
    add_text(
        sl,
        Inches(0.7),
        Inches(0.42),
        Inches(12),
        Inches(0.28),
        [("INTERNAL CONFIDENTIAL  ·  C-SUITE DECISION BRIEFING  ·  27 AUGUST 2026", 11, GOLD, True)],
    )
    add_text(
        sl,
        Inches(0.7),
        Inches(0.95),
        Inches(12.1),
        Inches(0.42),
        [("IMDS AGENTIC WORKFLOW", 14, GOLD, True)],
    )
    add_text(
        sl,
        Inches(0.7),
        Inches(1.35),
        Inches(12.1),
        Inches(1.15),
        [("Own the IMDS desk.\nSame-day accept. Zero PPAP delay.", 32, WHITE, True)],
    )
    add_text(
        sl,
        Inches(0.7),
        Inches(2.65),
        Inches(11.9),
        Inches(0.85),
        [
            (
                "50 MDS land in our inbox every working day. ~5,000 sit open — about 100 days of intake if nothing new arrived. Every one is a production-authorization document, not paperwork.",
                16,
                RGBColor(0xC5, 0xD0, 0xDC),
                False,
            )
        ],
    )
    meta = [
        ("Presented by", "Wong  ·  Kam Yuen Wong\nSupplier Quality Director / Data Scientist"),
        ("Company", "Johnson Electric\nInternational Limited"),
        ("Audience", "C-suite  ·  VP / GM\nOps, Supply Chain, Quality, Finance"),
        ("Today’s date", "Thursday\n27 August 2026"),
    ]
    x = Inches(0.7)
    for a, b in meta:
        card(sl, x, Inches(3.65), Inches(2.95), Inches(1.55), NAVY_MID, None)
        add_text(sl, x + Inches(0.14), Inches(3.74), Inches(2.67), Inches(0.24), [(a.upper(), 10, GOLD, True)])
        add_text(sl, x + Inches(0.14), Inches(4.02), Inches(2.67), Inches(1.05), [(b, 13, WHITE, False)])
        x += Inches(3.1)
    add_text(sl, Inches(0.7), Inches(5.40), Inches(12), Inches(0.28), [("TODAY’S ASK", 11, GOLD, True)])
    add_text(
        sl,
        Inches(0.7),
        Inches(5.70),
        Inches(12),
        Inches(1.15),
        [
            (
                "Invest in Agentic AI Today. Approve the 20-MDS live pilot and the internal budget line (implementation, training, maintenance). Name Supplier Quality as accountable owner, with executive sponsorship and a governance halt.",
                15,
                WHITE,
                False,
            )
        ],
    )
    notes(
        sl,
        _speech_goal(
            "1",
            "Do not walk a long agenda. Read the statistic once, then the ask once. Names in the room: Operations cares about PPAP, Supply Chain about the 5,000 open, Quality about wrongful-accept, Finance about payback. You are Wong, Supplier Quality Director / data scientist at Johnson Electric International Limited. This is a decision briefing, not a training class. After the title, the storyboard is the 60-second map of the next 20 minutes.",
        ),
    )


def s02_storyboard_a(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, WHITE)
    header(
        sl,
        "Storyboard  ·  presenter run sheet",
        "Eight beats. This is how we run the room.",
        "Use as a 60-second agenda, or skip and jump to Pain Points. Spoken hooks are on every notes page.",
    )
    rows = [["#", "Beat", "Slide intent", "Spoken hook", "Goal"]]
    for b in STORYBOARD[:4]:
        rows.append([b["n"], b["beat"], b["intent"], b["hook"], b["goal"]])
    table_shape = sl.shapes.add_table(len(rows), 5, Inches(0.40), Inches(1.42), Inches(12.50), Inches(5.45))
    table = table_shape.table
    widths = [0.55, 2.05, 3.15, 3.85, 2.90]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            if r == 0:
                set_cell(cell, val, 11, WHITE, True, NAVY, PP_ALIGN.CENTER)
            elif c == 0:
                set_cell(cell, val, 14, WHITE, True, TEAL, PP_ALIGN.CENTER)
            elif c == 1:
                set_cell(cell, val, 12, NAVY, True, LIGHT_GOLD if r % 2 else OFF)
            else:
                set_cell(cell, val, 11, SLATE, False, WHITE if r % 2 else OFF)
    footer(sl, 2)
    notes(
        sl,
        "STORYBOARD (beats 1–4). This slide is mandatory so a presenter can run the room. If the VP is already leaning in, spend 20 seconds: “Eight beats — pain, solution, impact, money, roadmap, proof, ask.” Then go. Do not teach IMDS UI here. Goal of the storyboard: keep you on the C-suite spine (ROI, risk, competitiveness, governance) and off engineering detail.",
    )


def s03_storyboard_b(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, WHITE)
    header(
        sl,
        "Storyboard  ·  presenter run sheet",
        "Beats 5–8: money, rollout, proof, the ask.",
        "Close on “Invest in Agentic AI Today” plus a specific pilot / budget approval — not a vague endorsement.",
    )
    rows = [["#", "Beat", "Slide intent", "Spoken hook", "Goal"]]
    for b in STORYBOARD[4:]:
        rows.append([b["n"], b["beat"], b["intent"], b["hook"], b["goal"]])
    table_shape = sl.shapes.add_table(len(rows), 5, Inches(0.40), Inches(1.42), Inches(12.50), Inches(5.45))
    table = table_shape.table
    widths = [0.55, 2.05, 3.15, 3.85, 2.90]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            if r == 0:
                set_cell(cell, val, 11, WHITE, True, NAVY, PP_ALIGN.CENTER)
            elif c == 0:
                set_cell(cell, val, 14, WHITE, True, COPPER, PP_ALIGN.CENTER)
            elif c == 1:
                set_cell(cell, val, 12, NAVY, True, LIGHT_GOLD if r % 2 else OFF)
            else:
                set_cell(cell, val, 11, SLATE, False, WHITE if r % 2 else OFF)
    footer(sl, 3)
    notes(
        sl,
        "STORYBOARD (beats 5–8). Timing target is 20 minutes including questions. Budget is internal-program hours, not a software RFP. Case studies are external / illustrative — say that out loud. CTA: “Invest in Agentic AI Today” and the 20-MDS pilot plus budget line. Executive sponsorship and stakeholder alignment — never military language.",
    )


def s04_pain(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, WHITE)
    header(
        sl,
        "02  ·  Pain points",
        "The IMDS desk is a production gate run as a manual inbox",
        "Frame this as ROI, launch risk, and competitiveness — not typing speed.",
    )
    kpis = [
        ("50", "MDS received / day", "Must be viewed before any accept or reject"),
        ("~5,000", "Open outstanding MDS", "~100 working days of intake if inbound froze"),
        ("8–60 min", "Per MDS, by hand", "Rec 001 + GM / VW / Ford overlays, node by node"),
        ("PPAP wait", "Cost of inaction", "Unaccepted MDS = open production authorization"),
    ]
    x = Inches(0.5)
    for v, lab, sub in kpis:
        kpi(sl, x, Inches(1.42), Inches(3.0), Inches(1.50), v, lab, sub, NAVY if v != "PPAP wait" else COPPER)
        x += Inches(3.15)

    pains = [
        ("ROI", "Specialist hours are spent on clicks — Check, accept, forward, propose, reject — not on supplier coaching or OEM-rule ownership. Hiring a parallel inbox does not encode three OEM rule packs."),
        ("Risk", "Passing IMDS Check is not OEM acceptance. IMDS 15.2 says further review may be required. Wrongful accept is a Quality event. A bouncing MDS to GM, VW, or Ford restarts PPAP."),
        ("Competitiveness", "We sit in two hops: Tier-2 to our customer and Tier-1 to the OEM. Dummy or unaccepted child MDS blocks the tree. Scorecards and launch dates do not wait for a 5,000-row backlog."),
        ("Capability", "Judgment sits in a few specialists’ heads. Leave or peak volume stalls launches. The scarce asset is the rule pack — not more people doing the same inbox."),
    ]
    y = Inches(3.12)
    x = Inches(0.5)
    colors = [TEAL, RED, COPPER, BLUE]
    for (t, b), col in zip(pains, colors):
        card(sl, x, y, Inches(3.0), Inches(3.85), WHITE, LINE)
        rect(sl, x, y, Inches(3.0), Inches(0.42), col)
        add_text(sl, x + Inches(0.14), y + Inches(0.08), Inches(2.72), Inches(0.28), [(t, 13, WHITE, True)])
        add_text(sl, x + Inches(0.14), y + Inches(0.55), Inches(2.72), Inches(3.15), [(b, 12, SLATE, False)])
        x += Inches(3.15)
    footer(sl, 4)
    notes(
        sl,
        _speech_goal(
            "2",
            "Four numbers, then one sentence per executive priority. Do not debate the 5,000 — it is why we are here. If challenged on “we already have IMDS,” confirm: classic browser only; no IMDS Plus, no Examiner, no Connect. HR in the room: this is not a reduction program. Quality: Rec 001 is the floor; GM / VW / Ford overlays are why first-pass dies. Operations: zero PPAP delay is the target, not a slogan.",
        ),
    )


def s05_solution(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, WHITE)
    header(
        sl,
        "03  ·  Proposed solution",
        "Agentic workflow: ingest → orchestrate → output",
        "A live inbox agent on our company IMDS account — not a nightly spreadsheet, not a vendor platform.",
    )
    stages = [
        (NAVY, "1  INGEST", "Received MDS rows from our IMDS inbox. Default batch: 20. One-button Colab run. Secrets stay in Colab (key icon): IMDS_USERNAME, IMDS_PASSWORD, OTP_SECRET."),
        (TEAL, "2  ORCHESTRATE", "Open each row. Run IMDS Check. Network-resume if the session drops. Leftover own-MDS sheets are closed cleanly so IDs never mismatch the next row."),
        (GREEN, "3  OUTPUT", "PASS → accept + forward + propose to recipients 9994 and 293798. FAIL → reject. Preferred contact Qu, Theresa; any real Supplier Data contact if she is missing."),
    ]
    x = Inches(0.45)
    for i, (color, title, body) in enumerate(stages):
        card(sl, x, Inches(1.38), Inches(3.85), Inches(3.35), WHITE, LINE)
        rect(sl, x, Inches(1.38), Inches(3.85), Inches(0.50), color)
        add_text(sl, x + Inches(0.16), Inches(1.46), Inches(3.53), Inches(0.36), [(title, 14, WHITE, True)])
        add_text(sl, x + Inches(0.16), Inches(2.02), Inches(3.53), Inches(2.50), [(body, 13, SLATE, False)])
        if i < 2:
            chevron(sl, x + Inches(3.78), Inches(2.85), Inches(0.28), Inches(0.36), GOLD)
        x += Inches(4.15)

    fork = [
        (GREEN, "PASS", "Accept the received MDS. Forward. Propose to company IDs 9994 and 293798. Contact: Qu, Theresa, with fallback to any real name in the dropdown."),
        (RED, "FAIL", "Reject with the Check result on the record. Do not forward. Do not propose. The specialist is not stuck in a click loop on a known-bad tree."),
    ]
    x = Inches(0.45)
    for color, title, body in fork:
        card(sl, x, Inches(4.90), Inches(6.15), Inches(2.05), WHITE, LINE)
        rect(sl, x, Inches(4.90), Inches(0.12), Inches(2.05), color)
        add_text(sl, x + Inches(0.28), Inches(5.05), Inches(2.2), Inches(0.32), [(title, 16, color, True)])
        add_text(sl, x + Inches(2.40), Inches(5.08), Inches(3.55), Inches(1.70), [(body, 13, SLATE, False)])
        x += Inches(6.30)
    footer(sl, 5)
    notes(
        sl,
        _speech_goal(
            "3",
            "Walk left to right: ingest, orchestrate, output. Then the fork. Measurable outcomes: a 20-MDS batch decided the same session; PASS rows land at 9994 and 293798; FAIL rows are rejected, not left aging. Position this as a strategic production-gate control, not “we wrote a script.” If asked about credentials: Colab Secrets only — never in git. If asked about contact: Theresa is preferred; any real Supplier Data contact is allowed so a missing name does not stall Propose.",
        ),
    )


def s06_live_agent(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, WHITE)
    header(
        sl,
        "03  ·  Proposed solution",
        "The live agent is already hardened for a governed pilot",
        "imds_agent_v2.py on our licensed IMDS users only. Reliability is the governance story.",
    )
    items = [
        ("One-button Colab", "Run from Colab_Start_Here. Credentials in Colab Secrets (key icon). Default 20 MDS. Leftover debug values of 3 or 10 are ignored so an old cell cannot shrink the pilot."),
        ("PASS path", "IMDS Check with 0 errors → accept, then forward, then propose. Recipients 9994 and 293798. Preferred contact Qu, Theresa."),
        ("FAIL path", "Check failures are rejected. The inbox does not accumulate known-bad MDS while specialists hunt clicks."),
        ("Network resume", "A dropped session waits and continues the 20-row loop instead of aborting the batch. Production gates cannot depend on a perfect network afternoon."),
        ("Same-MDS tabs", "“Do you want to save your changes?” is Yes when switching Supplier Data / Recipient Data on the same forwarded MDS so Propose finishes on the new own-MDS ID."),
        ("Leaving leftovers", "The same prompt is No when leaving a leftover own-MDS sheet to search the next received row — Yes would keep the leftover ID and mismatch every later row."),
    ]
    y = Inches(1.32)
    x0 = Inches(0.5)
    for i, (t, b) in enumerate(items):
        col = i % 2
        row = i // 2
        x = x0 + Inches(6.35) * col
        yy = y + Inches(1.82) * row
        card(sl, x, yy, Inches(6.15), Inches(1.68), WHITE, LINE)
        oval(sl, x + Inches(0.16), yy + Inches(0.18), Inches(0.36), Inches(0.36), TEAL)
        add_text(sl, x + Inches(0.16), yy + Inches(0.22), Inches(0.36), Inches(0.30), [(str(i + 1), 12, WHITE, True, PP_ALIGN.CENTER)])
        add_text(sl, x + Inches(0.64), yy + Inches(0.18), Inches(5.30), Inches(0.32), [(t, 15, NAVY, True)])
        add_text(sl, x + Inches(0.64), yy + Inches(0.54), Inches(5.30), Inches(1.00), [(b, 12, SLATE, False)])
    footer(sl, 6)
    notes(
        sl,
        _speech_goal(
            "3",
            "This is still Beat 3. C-suite translation: the agent is not a lab toy. Network resume, leftover-sheet handling, and the save-changes rule are how we keep a 20-row pilot from corrupting IDs. Own-account only — our credentials, our inbox. No other company’s IMDS data. Emergency halt (governance halt) is the control if Quality wants auto-accept paused. Do not open a DXC commercial discussion unless asked.",
        ),
    )


def s07_impact(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, WHITE)
    header(
        sl,
        "04  ·  Business impact",
        "Manual inbox versus the agentic workflow",
        "Transformation is hours, launch risk, and a decision log — not a software logo.",
    )
    rows = [
        ["Work", "Before — manual desk", "After — agentic workflow"],
        ["Inbox", "50/day viewed by hand; ~5,000 aging", "20-MDS governed batch; same-session decisions"],
        ["Check", "Specialist walks every node, every overlay", "Agent runs IMDS Check; PASS/FAIL is explicit"],
        ["PASS path", "Accept, forward, propose as separate click chains", "Accept + forward + propose to 9994 and 293798"],
        ["FAIL path", "Rejects wait; reasons vary by who is on shift", "Reject on Check failure; reason on the record"],
        ["Contact / recipients", "Theresa / org IDs re-typed; easy to miss", "Preferred contact with fallback; both recipients every PASS"],
        ["Disruption", "Network drop or leftover sheet aborts the afternoon", "Resume after drop; leftover sheets closed with No"],
        ["PPAP / scorecard", "Open MDS = open launch authorization", "Same-day greens; reds returned to the supplier the same day"],
        ["Audit trail", "Tribal knowledge in specialists’ heads", "Excel check_summary with Action Result per row"],
    ]
    table_shape = sl.shapes.add_table(len(rows), 3, Inches(0.40), Inches(1.32), Inches(12.50), Inches(5.55))
    table = table_shape.table
    table.columns[0].width = Inches(2.15)
    table.columns[1].width = Inches(5.15)
    table.columns[2].width = Inches(5.20)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            if r == 0:
                fill = NAVY if c != 2 else TEAL
                set_cell(cell, val, 12, WHITE, True, fill, PP_ALIGN.CENTER)
            elif c == 0:
                set_cell(cell, val, 12, NAVY, True, LIGHT_GOLD if r % 2 else OFF)
            elif c == 1:
                set_cell(cell, val, 12, SLATE, False, LIGHT_RED if r % 2 else OFF)
            else:
                set_cell(cell, val, 12, SLATE, False, LIGHT_TEAL if r % 2 else WHITE)
    footer(sl, 7)
    notes(
        sl,
        _speech_goal(
            "4",
            "Read three rows only if short on time: Inbox, PASS path, PPAP. Efficiency: clicks leave the specialist. Cost: hours and launch delay. New capacity (not a claimed new revenue stream we have not booked): specialists become OEM-rule owners and supplier coaches — the work that actually protects customer scorecards. Do not invent a Johnson Electric dollar savings figure. The 50/day and 5,000 open are our desk facts; keep them.",
        ),
    )


def s08_budget(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, WHITE)
    header(
        sl,
        "05  ·  Budget & ROI",
        "Internal program hours. Zero software license. Payback inside a quarter.",
        "Planning model for this briefing — not booked Johnson Electric P&L. No vendor invoice.",
    )
    kpis = [
        ("$0", "Software licenses", "Classic IMDS browser + Colab Secrets. No Plus, no Examiner, no integrator."),
        ("~60 h", "Implementation", "Govern the live agent: logging, sample audit, halt, recipient/contact policy."),
        ("~16 h", "Training", "Specialists: exception handling, audit sampling, emergency halt drill."),
        ("~4 h/mo", "Maintenance", "Rule drift, recipients, contact fallback, secret hygiene, weekly sample audit."),
    ]
    x = Inches(0.45)
    for v, lab, sub in kpis:
        kpi(sl, x, Inches(1.38), Inches(3.05), Inches(1.55), v, lab, sub, TEAL if v == "$0" else NAVY)
        x += Inches(3.20)

    rows = [
        ["ROI lever", "How we count it", "Why it pays back"],
        [
            "Hours on Check / accept / forward / propose",
            "15 min median × 50 MDS/day ≈ 12.5 specialist-hours/day (desk model)",
            "Even ~25% of inbound automated returns 3+ hours/day. 60 implementation hours recover in weeks.",
        ],
        [
            "IMDS non-compliance / OEM scorecards",
            "GM, VW, Ford completeness and bounce rate on our MDS",
            "A scorecard miss is commercial risk. Logged PASS/FAIL is the control Finance and Quality can see.",
        ],
        [
            "Delayed PPAP / launch hold",
            "Unaccepted or bouncing MDS blocks production authorization",
            "One delayed PPAP dwarfs the pilot cost. Same-day reject also shortens supplier rework.",
        ],
        [
            "Payback period",
            "Target: inside 90 days of pilot start",
            "Hours recovered + avoided launch delay. Revisit after the 20-MDS pilot with measured minutes/MDS.",
        ],
    ]
    table_shape = sl.shapes.add_table(len(rows), 3, Inches(0.40), Inches(3.12), Inches(12.50), Inches(3.75))
    table = table_shape.table
    table.columns[0].width = Inches(3.30)
    table.columns[1].width = Inches(4.70)
    table.columns[2].width = Inches(4.50)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            if r == 0:
                set_cell(cell, val, 12, WHITE, True, NAVY, PP_ALIGN.CENTER)
            elif c == 0:
                set_cell(cell, val, 12, NAVY, True, LIGHT_GOLD if r % 2 else OFF)
            else:
                set_cell(cell, val, 11, SLATE, False, WHITE if r % 2 else OFF)
    footer(sl, 8)
    notes(
        sl,
        _speech_goal(
            "5",
            "Emphasize: this is own effort, not a purchase. Software is $0. Implementation is operationalizing an agent that already runs. Training is role shift, not a classroom tour. Maintenance is hours per month, not a platform contract. Payback: say “inside one quarter on hours alone” and “one PPAP delay exceeds the pilot.” If Finance wants a dollar ROI, offer to price the 12.5 hours/day at our loaded specialist rate after the meeting — do not invent a figure on the slide. Risk mitigation: emergency halt, sample audit, own-account only.",
        ),
    )


def s09_roadmap(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, WHITE)
    header(
        sl,
        "06  ·  Implementation roadmap",
        "Pilot → Scale → Optimize. Controlled rollout, named governance.",
        "We do not flip auto-accept on the 5,000-row stock on day one.",
    )
    phases = [
        (TEAL, "PILOT", "Days 1–30", "20 MDS live batches on our company account. Log every PASS/FAIL. Sample-audit greens. Freeze recipient IDs 9994 / 293798 and the Theresa-plus-fallback contact rule. Emergency halt stays in Supplier Quality’s hands. Success: the 20-row loop completes without leftover-ID mismatch."),
        (BLUE, "SCALE", "Days 31–60", "Raise daily burn on NEW inbound. Auto-reject FAIL patterns with structured text. Shadow-score remaining greens beside the specialist. Operations uses PPAP clock as the success metric. Supply chain chases dummy children the agent cannot invent."),
        (NAVY, "OPTIMIZE", "Days 61–90+", "Auto-accept greens with weekly audit. Start a measured burn-down of the 5,000. Encode GM / VW / Ford overlays as Examiner-equivalent packs on top of Check. Role shift: specialists own rule packs and exceptions — not inbox volume."),
    ]
    x = Inches(0.45)
    for color, title, when, body in phases:
        card(sl, x, Inches(1.38), Inches(4.05), Inches(4.05), WHITE, LINE)
        rect(sl, x, Inches(1.38), Inches(4.05), Inches(0.85), color)
        add_text(sl, x + Inches(0.18), Inches(1.46), Inches(3.70), Inches(0.32), [(title, 16, WHITE, True)])
        add_text(sl, x + Inches(0.18), Inches(1.80), Inches(3.70), Inches(0.32), [(when, 12, GOLD, False)])
        add_text(sl, x + Inches(0.18), Inches(2.38), Inches(3.70), Inches(2.85), [(body, 13, SLATE, False)])
        x += Inches(4.20)

    card(sl, Inches(0.45), Inches(5.58), Inches(12.40), Inches(1.38), LIGHT_GOLD, GOLD)
    add_text(sl, Inches(0.65), Inches(5.68), Inches(12.00), Inches(0.28), [("GOVERNANCE — non-negotiable", 12, NAVY, True)])
    add_text(
        sl,
        Inches(0.65),
        Inches(6.00),
        Inches(12.00),
        Inches(0.82),
        [
            (
                "Accountable owner: Supplier Quality Director. Executive sponsorship from VP/GM so the pilot is not re-argued every week. Own-account only. Sample audit of PASS accepts. Emergency halt on auto-accept / auto-reject in one action. Human-mandatory list: novel chemistry, OEM derogation, dispute. Stakeholder alignment: Quality (audit), Operations (PPAP clock), Supply chain (dummy-child follow-up), HR (role redesign, not a headcount cut).",
                13,
                SLATE,
                False,
            )
        ],
    )
    footer(sl, 9)
    notes(
        sl,
        _speech_goal(
            "6",
            "Stress controlled rollout. Pilot is 20 MDS, not the 5,000. Scale is new inbound. Optimize is when auto-accept is earned. Governance: executive sponsorship, emergency halt, sample audit, own-account only. If Legal flinches at auto-accept, keep it in Optimize and run FAIL-reject + logged PASS in Pilot. HR: say the desk does not disappear — the job becomes OEM-rule ownership.",
        ),
    )


def s10_cases(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, WHITE)
    header(
        sl,
        "07  ·  Case studies",
        "External proof the category works. Not Johnson Electric results.",
        "Vendor-reported and industry-analog. We keep OEM knowledge in-house; we do not rent the desk.",
    )
    card(sl, Inches(0.45), Inches(1.32), Inches(12.40), Inches(0.52), LIGHT_GOLD, GOLD)
    add_text(
        sl,
        Inches(0.65),
        Inches(1.40),
        Inches(12.00),
        Inches(0.36),
        [("ILLUSTRATIVE / EXTERNAL — do not quote these as Johnson Electric metrics", 14, NAVY, True)],
    )
    cases = [
        (
            "Automotive IMDS platforms",
            "DXC IMDS Plus / Inbox Automation; iPoint; APA MDS Xpress",
            "Paid Examiner profiles and auto-accept if no errors. Proves OEMs already accept machine-paced inbox work. We have none of it, and we are not opening an RFP in v1.",
        ),
        (
            "Predco — vendor-reported, 12 Tier-1s",
            "Automotive supply-chain compliance automation",
            "−89% rejects, −72% cycle time, ~18 days PPAP recovery, 1.8% reject rate. Evidence the problem is solvable. Not a purchase recommendation. Not our ROI.",
        ),
        (
            "Certivo — vendor-reported",
            "Agentic / automated MDS preparation",
            "First-pass reject 30–40% → <5%; cycle 4–6 weeks → ~4 hours; ~80% less prep labor. Productivity and customer-readiness analog. Not booked JE revenue lift.",
        ),
    ]
    x = Inches(0.45)
    for title, sub, body in cases:
        card(sl, x, Inches(2.02), Inches(4.05), Inches(3.55), WHITE, LINE)
        rect(sl, x, Inches(2.02), Inches(4.05), Inches(0.85), TEAL)
        add_text(sl, x + Inches(0.16), Inches(2.10), Inches(3.73), Inches(0.40), [(title, 13, WHITE, True)])
        add_text(sl, x + Inches(0.16), Inches(2.50), Inches(3.73), Inches(0.30), [(sub, 11, LIGHT_GOLD, False)])
        add_text(sl, x + Inches(0.16), Inches(3.02), Inches(3.73), Inches(2.35), [(body, 13, SLATE, False)])
        x += Inches(4.20)

    add_text(
        sl,
        Inches(0.50),
        Inches(5.72),
        Inches(12.30),
        Inches(1.20),
        [
            (
                "What we take from them: score, cite, structured reject, outbound pre-flight. What we do not do: outsource GM / VW / Ford overlays or our 5,000-row desk to a vendor that does not know our tree conventions. Johnson Electric metrics stay the ones on the pain and impact slides (50/day, ~5,000 open, 20-MDS live agent). Customer-satisfaction and revenue-lift figures above are theirs, not ours.",
                14,
                NAVY,
                False,
            )
        ],
    )
    footer(sl, 10)
    notes(
        sl,
        _speech_goal(
            "7",
            "Say “external / illustrative” before any number. Predco and Certivo are vendor-reported; they reduce perceived risk, they are not our scorecard. If someone says “just buy Plus,” answer: Plus is a paid Examiner; it does not encode GMW3059 / VW 91101 / Ford RSMS for us, and we have no software budget. Build captures the rule packs we already apply by hand. Productivity analog is enough; do not claim a JE customer-satisfaction or revenue-lift number we have not measured.",
        ),
    )


def s11_cta(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, NAVY)
    rect(sl, 0, 0, Inches(0.18), H, GOLD)
    add_text(
        sl,
        Inches(0.7),
        Inches(0.40),
        Inches(12.0),
        Inches(0.28),
        [("08  ·  CALL TO ACTION  ·  27 AUGUST 2026", 12, GOLD, True)],
    )
    add_text(
        sl,
        Inches(0.7),
        Inches(0.90),
        Inches(12.1),
        Inches(1.10),
        [("Invest in Agentic AI Today.", 36, WHITE, True)],
    )
    add_text(
        sl,
        Inches(0.7),
        Inches(2.10),
        Inches(12.0),
        Inches(0.55),
        [
            (
                "Approve the 20-MDS live pilot and the internal budget line this week.",
                20,
                GOLD,
                True,
            )
        ],
    )
    asks = [
        ("1", "Pilot", "Authorize live runs of 20 received MDS on our company IMDS account (Colab one-button, Secrets in 🔑)."),
        ("2", "Budget", "Allocate the internal line: ~60 h implementation, ~16 h training, ~4 h/month maintenance. $0 licenses."),
        ("3", "Owner", "Supplier Quality Director (Wong) owns rule packs, emergency halt, 90-day milestones, and the 5,000 burn-down."),
        ("4", "Sponsorship", "VP/GM executive sponsorship and stakeholder alignment so the pilot is the house approach — not a weekly re-litigation of buy vs hire."),
    ]
    y = Inches(2.80)
    for n, t, b in asks:
        card(sl, Inches(0.7), y, Inches(11.9), Inches(0.85), NAVY_MID, None)
        oval(sl, Inches(0.88), y + Inches(0.22), Inches(0.42), Inches(0.42), GOLD)
        add_text(sl, Inches(0.88), y + Inches(0.28), Inches(0.42), Inches(0.32), [(n, 14, NAVY, True, PP_ALIGN.CENTER)])
        add_text(sl, Inches(1.50), y + Inches(0.12), Inches(1.70), Inches(0.60), [(t, 16, WHITE, True)])
        add_text(sl, Inches(3.30), y + Inches(0.16), Inches(9.05), Inches(0.58), [(b, 14, RGBColor(0xD5, 0xDE, 0xE8), False)])
        y += Inches(0.95)
    notes(
        sl,
        _speech_goal(
            "8",
            "Read the bold line exactly: “Invest in Agentic AI Today.” Then the specific ask: approve the 20-MDS pilot and the budget line. Walk the four boxes. Stop. Leave silence so VP/GM can endorse. If they ask budget, say own effort only — 60 / 16 / 4 hours. If they ask headcount, say role redesign not reduction. If they ask buy vs build, say build; vendors are evidence. Do not add a new idea after the last line.",
        ),
    )


def s12_appendix(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, WHITE)
    header(
        sl,
        "Appendix  ·  Q&A only",
        "Live-agent facts and glossary — skip in the 20-minute slot",
        "Validate OEM overlays against current GM / VW / Ford IMDS portals before production use.",
    )
    glossary = [
        ("Executive sponsorship", "VP/GM publicly backs the pilot so buy-vs-hire is not reopened weekly."),
        ("Emergency halt", "One action stops auto-accept / auto-reject if a pack or Check rule drifts."),
        ("Own-account only", "Our IMDS users, our company inbox/outbox. No other company’s data."),
        ("PASS", "IMDS Check with 0 errors → accept + forward + propose to 9994, 293798."),
        ("FAIL", "Check failure → reject. No forward, no propose."),
        ("Preferred contact", "Qu, Theresa. Fallback: any real Supplier Data name (not blank)."),
        ("Network resume", "After a drop, wait and continue the 20-row loop instead of aborting."),
        ("Save-changes Yes", "Same forwarded MDS, switching Supplier / Recipient tabs."),
        ("Save-changes No", "Leaving a leftover own-MDS sheet to open the next received row."),
        ("Dummy child", "Placeholder node instead of an accepted sub-supplier MDS."),
        ("Pre-flight", "Run the OEM overlay on outbound MDS before Propose."),
        ("OEM-rule owner", "Specialist owns GM/VW/Ford packs and exceptions, not inbox volume."),
    ]
    rows = [["Term", "Meaning", "Term", "Meaning"]]
    for i in range(0, len(glossary), 2):
        a, b = glossary[i]
        c, d = glossary[i + 1] if i + 1 < len(glossary) else ("", "")
        rows.append([a, b, c, d])
    table_shape = sl.shapes.add_table(len(rows), 4, Inches(0.40), Inches(1.28), Inches(12.50), Inches(5.55))
    table = table_shape.table
    widths = [2.15, 4.10, 2.15, 4.10]
    for i, w in enumerate(widths):
        table.columns[i].width = Inches(w)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            if r == 0:
                set_cell(cell, val, 11, WHITE, True, NAVY, PP_ALIGN.CENTER)
            elif c in (0, 2):
                set_cell(cell, val, 11, NAVY, True, LIGHT_GOLD if r % 2 else OFF)
            else:
                set_cell(cell, val, 11, SLATE, False, WHITE if r % 2 else OFF)
    footer(sl, 12)
    notes(
        sl,
        "APPENDIX — do not present in the 20-minute slot unless someone asks for definitions. Directed-buy example if asked: OEM nominates a sub-supplier; they Propose to the OEM and to us. Dummy-child follow-up: Supply chain emails sub-suppliers with missing accepted child MDS. Orphan decision: accept/reject with no rule-pack version on record. Sources: public.mdsystem.com OEM guides, GMW3059, VW 91101, Ford RSMS. Live agent: imds_agent_v2.py; Colab Secrets IMDS_USERNAME, IMDS_PASSWORD, OTP_SECRET.",
    )


def all_deck_text(prs: Presentation) -> str:
    chunks: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        chunks.append(cell.text_frame.text)
        notes_slide = slide.notes_slide
        chunks.append(notes_slide.notes_text_frame.text)
    return "\n".join(chunks)


def assert_no_forbidden_jargon(prs: Presentation) -> None:
    blob = all_deck_text(prs).lower()
    hits = [term for term in FORBIDDEN_JARGON if term in blob]
    if hits:
        raise SystemExit(f"Forbidden jargon found in deck: {hits}")


def export_docs(pptx_path: Path) -> None:
    """Copy pptx into docs/ and render PDF + PNG slides when LibreOffice is available."""
    DOCS_PPTX.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(pptx_path, DOCS_PPTX)
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        print("LibreOffice not found — skipped PDF/PNG export", file=sys.stderr)
        return
    outdir = ROOT / "docs"
    env = os.environ.copy()
    env.setdefault("HOME", "/tmp")
    subprocess.run(
        [soffice, "--headless", "--norestore", "--convert-to", "pdf", "--outdir", str(outdir), str(pptx_path)],
        check=True,
        env=env,
    )
    pdf = DOCS_PDF
    if not pdf.exists():
        print("PDF was not produced", file=sys.stderr)
        return
    SLIDES_DIR.mkdir(parents=True, exist_ok=True)
    for old in SLIDES_DIR.glob("slide-*.png"):
        old.unlink()
    prefix = SLIDES_DIR / "slide"
    subprocess.run(["pdftoppm", "-png", "-r", "140", str(pdf), str(prefix)], check=True)
    # pdftoppm writes slide-1.png; normalize to slide-01.png
    produced = sorted(SLIDES_DIR.glob("slide-*.png"), key=lambda p: int("".join(ch for ch in p.stem if ch.isdigit()) or "0"))
    for i, path in enumerate(produced, 1):
        dest = SLIDES_DIR / f"slide-{i:02d}.png"
        if path != dest:
            if dest.exists():
                dest.unlink()
            path.rename(dest)
    print(f"Exported {DOCS_PDF} and {len(list(SLIDES_DIR.glob('slide-*.png')))} PNGs")


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    builders = [
        s01_title,
        s02_storyboard_a,
        s03_storyboard_b,
        s04_pain,
        s05_solution,
        s06_live_agent,
        s07_impact,
        s08_budget,
        s09_roadmap,
        s10_cases,
        s11_cta,
        s12_appendix,
    ]
    global TOTAL
    TOTAL = len(builders)
    for fn in builders:
        fn(prs)
    assert_no_forbidden_jargon(prs)
    return prs


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--export-docs", action="store_true", help="Also write docs/ pptx, PDF, and PNG slides")
    args = parser.parse_args(argv)
    prs = build()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT} ({TOTAL} slides)")
    if args.export_docs:
        export_docs(OUT)
        shutil.copy2(OUT, DOCS_PPTX)
    else:
        DOCS_PPTX.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(OUT, DOCS_PPTX)
        print(f"Copied {DOCS_PPTX}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
