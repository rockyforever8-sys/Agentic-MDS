#!/usr/bin/env python3
"""Guard the C-suite IMDS briefing: storyboard spine, no military jargon."""

from __future__ import annotations

import sys
import tempfile
import unittest
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
if str(ROOT) not in sys.path:
    sys.path.insert(0, str(ROOT))

import importlib.util

from pptx import Presentation

_spec = importlib.util.spec_from_file_location(
    "build_imds_presentation",
    ROOT / "scripts" / "build_imds_presentation.py",
)
briefing = importlib.util.module_from_spec(_spec)
assert _spec.loader is not None
_spec.loader.exec_module(briefing)

FORBIDDEN_JARGON = briefing.FORBIDDEN_JARGON
STORYBOARD = briefing.STORYBOARD
all_deck_text = briefing.all_deck_text
build = briefing.build


def _notes_text(slide) -> str:
    return slide.notes_slide.notes_text_frame.text.strip()


class PresentationBriefingTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls):
        cls.prs = build()
        cls.blob = all_deck_text(cls.prs)
        cls.lower = cls.blob.lower()

    def test_storyboard_has_eight_beats(self):
        self.assertEqual(len(STORYBOARD), 8)
        names = [row["beat"] for row in STORYBOARD]
        self.assertEqual(
            names,
            [
                "Executive Opening",
                "Pain Points",
                "Proposed Solution",
                "Business Impact",
                "Budget & ROI",
                "Implementation Roadmap",
                "Case Studies",
                "Call to Action",
            ],
        )
        for row in STORYBOARD:
            self.assertTrue(row["intent"])
            self.assertTrue(row["hook"])
            self.assertTrue(row["goal"])

    def test_storyboard_slides_map_beats(self):
        combined = self.blob
        for row in STORYBOARD:
            self.assertIn(row["beat"], combined)
            self.assertIn(row["intent"], combined)
            self.assertIn(row["hook"], combined)
            self.assertIn(row["goal"], combined)

    def test_slide_count_is_executive_length(self):
        self.assertEqual(len(self.prs.slides), 12)
        self.assertLessEqual(len(self.prs.slides), 16)

    def test_every_slide_has_speaker_notes(self):
        for i, slide in enumerate(self.prs.slides, 1):
            text = _notes_text(slide)
            self.assertGreater(len(text), 40, f"slide {i} notes too short")

    def test_title_has_presenter_project_and_date(self):
        self.assertIn("IMDS AGENTIC WORKFLOW", self.blob)
        self.assertIn("Kam Yuen Wong", self.blob)
        self.assertIn("Johnson Electric", self.blob)
        self.assertIn("27 August 2026", self.blob)
        self.assertIn("Supplier Quality Director", self.blob)

    def test_cta_is_explicit(self):
        self.assertIn("Invest in Agentic AI Today.", self.blob)
        self.assertIn("20-MDS", self.blob)
        self.assertIn("budget line", self.lower)

    def test_live_agent_facts(self):
        for token in (
            "20",
            "9994",
            "293798",
            "Qu, Theresa",
            "IMDS_USERNAME",
            "IMDS_PASSWORD",
            "OTP_SECRET",
            "Colab",
            "accept",
            "forward",
            "propose",
            "reject",
            "Network resume",
        ):
            self.assertIn(token, self.blob)

    def test_je_desk_metrics_kept(self):
        self.assertIn("5,000", self.blob)
        self.assertIn("50", self.blob)
        self.assertIn("PPAP", self.blob)
        self.assertIn("GM", self.blob)

    def test_case_studies_are_labeled_external(self):
        self.assertIn("ILLUSTRATIVE / EXTERNAL", self.blob)
        self.assertIn("vendor-reported", self.lower)
        self.assertIn("Not Johnson Electric results", self.blob)

    def test_budget_has_three_cost_lines(self):
        self.assertIn("Implementation", self.blob)
        self.assertIn("Training", self.blob)
        self.assertIn("Maintenance", self.blob)
        self.assertIn("Payback", self.blob)

    def test_roadmap_phases(self):
        self.assertIn("PILOT", self.blob)
        self.assertIn("SCALE", self.blob)
        self.assertIn("OPTIMIZE", self.blob)

    def test_no_military_jargon(self):
        for term in FORBIDDEN_JARGON:
            self.assertNotIn(term, self.lower)
        for term in ("air cover", "cover fire", "battlespace", "war room", "kill chain", "kill switch"):
            self.assertNotIn(term, self.lower)

    def test_executive_language_replaces_air_cover(self):
        self.assertIn("executive sponsorship", self.lower)
        self.assertIn("emergency halt", self.lower)
        self.assertIn("stakeholder alignment", self.lower)

    def test_saves_without_error(self):
        with tempfile.TemporaryDirectory() as tmp:
            path = Path(tmp) / "briefing.pptx"
            self.prs.save(str(path))
            again = Presentation(str(path))
            self.assertEqual(len(again.slides), 12)


if __name__ == "__main__":
    unittest.main()
